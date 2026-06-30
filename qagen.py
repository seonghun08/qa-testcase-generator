#!/usr/bin/env python3
"""
QA 테스트 케이스 문서 자동 생성 CLI

사용법:
  python qagen.py INFOSYS-4754 --ppt ./decrypted.pptx
  python qagen.py INFOSYS-4754 --jira-text ./jira.txt
  python qagen.py INFOSYS-4754 --ppt ./decrypted.pptx --jira-text ./jira.txt
  python qagen.py INFOSYS-4754 --interactive
"""
import argparse
import json
import os
import sys
import textwrap

from dotenv import load_dotenv

from template_builder import (
    build_workbook,
    save_workbook,
    LOG_TEMPLATE,
    DB_TEMPLATES,
    SYSTEM_TYPES,
)


# ============================================================
# PPT 텍스트 추출
# ============================================================

def extract_ppt_text(ppt_path):
    """python-pptx로 PPT에서 슬라이드별 텍스트를 추출한다."""
    from pptx import Presentation

    if not os.path.isfile(ppt_path):
        print(f"[오류] PPT 파일을 찾을 수 없습니다: {ppt_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(ppt_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        if texts:
            slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


# ============================================================
# 텍스트 파일 읽기
# ============================================================

def read_text_file(path):
    """텍스트 파일을 읽어 반환한다."""
    if not os.path.isfile(path):
        print(f"[오류] 파일을 찾을 수 없습니다: {path}", file=sys.stderr)
        sys.exit(1)
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()


# ============================================================
# Claude API 호출
# ============================================================

SYSTEM_PROMPT = textwrap.dedent("""\
    당신은 PG사 백오피스 QA 테스트 케이스 작성 전문가입니다.
    개발 요청서(PPT 또는 Jira 티켓)를 분석하여 사내 표준 QA 테스트 케이스를 생성합니다.

    ## 출력 형식
    반드시 아래 JSON 형식으로만 응답하세요. 설명이나 마크다운 없이 순수 JSON만 출력하세요.

    ```json
    {
      "title": "문서 제목 (예: 가맹점 정산 내역 관리 화면 개발 요청)",
      "system_type": "BO 또는 SHOP 또는 OSS",
      "test_cases": [
        {
          "id": "TB_01_001",
          "precond": "사전조건 (예: 내부직원 계정 로그인\\n\\nPG 정산 관리 > 가맹점 정산 내역)",
          "step": "조회|등록|수정|삭제",
          "category": "구분 (예: 정산 내역 목록 조회)",
          "item": "항목 (예: 검색 조건 별 목록 조회)",
          "procedure": "수행절차 (번호 매기기)",
          "expected": "기대결과 (번호 매기기)",
          "log_db_type": "select|insert|update"
        }
      ]
    }
    ```

    ## TestCaseID 규칙
    - 형식: TB_{기능그룹번호}_{순번} (예: TB_01_001)
    - 화면 단위로 그룹 분리 (메인 화면=01, 상세 팝업=02 등)
    - 한 화면 안에서는 조회→등록→수정→삭제 순
    - 순번은 그룹마다 001부터 시작

    ## 사전조건 형식
    - "{계정유형} 계정 로그인\\n\\n{메뉴경로}"
    - 시스템 타입별 계정: BO→내부직원, SHOP→가맹점관리자, OSS→운영지원담당자

    ## 테스트 케이스 작성 지침
    - 수행절차: 사용자 관점의 구체적 조작 단계 (번호+하위번호 형식)
    - 기대결과: 각 단계에 대응하는 검증 포인트
    - log_db_type: 해당 케이스의 DB 조작 유형
      - select: 조회 기능
      - insert: 등록/생성 기능
      - update: 수정/삭제/상태변경 기능

    ## 중요 규칙
    - 분량: 케이스 개수를 미리 정하지 않는다. 'QA 주제'가 지정되면 그 주제에만 한정해 주제당 1건(많아야 2건)으로 작성하고, 나열되지 않은 기능/화면은 만들지 않는다
      - 'QA 주제'가 없으면 요청서에서 핵심 기능을 뽑아 화면/기능 단위로 크게 묶어 간결하게 작성 (세부 단계까지 쪼개지 않음)
      - 한 주제(케이스) 안에서 세부 동작·상태별 분기·검증 포인트는 기대결과에 하위 번호로 나열
      - 확인 팝업/필수 입력 검증/정상 처리/실패 처리/이력 적재는 별도 케이스로 분리하지 말고 한 케이스로 통합
    - 정상 케이스뿐만 아니라 경계값, 에러 케이스도 포함
    - 단, 권한 체크와 상태 변경 정책(변경 불가) 등 핵심 에러 케이스는 통합하되 누락 금지
    - 실제 현업에서 쓸 수 있는 수준의 상세함 유지
""")

USER_PROMPT_TEMPLATE = textwrap.dedent("""\
    아래 개발 요청서를 분석하여 QA 테스트 케이스를 생성해주세요.
    Jira 티켓 번호: {jira_num}

    ---
    {input_text}
    ---
    {qa_topics}
    위 요청서를 기반으로 테스트 케이스 JSON을 생성해주세요.
""")


def call_claude_api(jira_num, input_text, api_key, qa_topics=""):
    """Claude API를 호출하여 테스트 케이스 JSON을 받아온다."""
    import anthropic

    client = anthropic.Anthropic(api_key=api_key)

    topics_section = ""
    if qa_topics.strip():
        topics_section = (
            "\n## QA 주제 (사용자 지정 테스트 범위 — 반드시 이 주제들에만 한정)\n"
            f"{qa_topics}\n"
            "- 위 주제만 작성하고 주제당 1건(많아야 2건)으로 통합, 나열되지 않은 기능/화면은 생성하지 않는다.\n"
        )

    user_prompt = USER_PROMPT_TEMPLATE.format(
        jira_num=jira_num, input_text=input_text, qa_topics=topics_section
    )

    print("[*] Claude API 호출 중...")
    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=8192,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": user_prompt}],
    )

    response_text = message.content[0].text
    print(f"[*] 응답 수신 완료 (토큰: 입력 {message.usage.input_tokens}, 출력 {message.usage.output_tokens})")

    return parse_response(response_text)


def parse_response(response_text):
    """Claude 응답에서 JSON을 파싱한다."""
    text = response_text.strip()

    # ```json ... ``` 블록이 있으면 그 안의 내용만 추출
    if '```json' in text:
        start = text.index('```json') + len('```json')
        end = text.index('```', start)
        text = text[start:end].strip()
    elif '```' in text:
        start = text.index('```') + 3
        end = text.index('```', start)
        text = text[start:end].strip()

    try:
        data = json.loads(text)
    except json.JSONDecodeError as e:
        print(f"[오류] JSON 파싱 실패: {e}", file=sys.stderr)
        print(f"[디버그] 원본 응답:\n{response_text[:500]}", file=sys.stderr)
        sys.exit(1)

    if 'test_cases' not in data:
        print("[오류] 응답에 test_cases 필드가 없습니다.", file=sys.stderr)
        sys.exit(1)

    return data


# ============================================================
# 인터랙티브 모드
# ============================================================

def interactive_input():
    """사용자로부터 직접 텍스트 입력을 받는다."""
    print("=" * 60)
    print("개발 요청 내용을 입력하세요.")
    print("입력 완료 후 빈 줄에서 Ctrl+Z (Windows) 또는 Ctrl+D (Mac/Linux)를 누르세요.")
    print("=" * 60)
    lines = []
    try:
        while True:
            line = input()
            lines.append(line)
    except EOFError:
        pass
    text = "\n".join(lines)
    if not text.strip():
        print("[오류] 입력된 내용이 없습니다.", file=sys.stderr)
        sys.exit(1)
    return text


# ============================================================
# 메인
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="QA 테스트 케이스 문서(xlsx) 자동 생성 CLI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
            사용 예시:
              python qagen.py INFOSYS-4754 --ppt ./decrypted.pptx
              python qagen.py INFOSYS-4754 --jira-text ./jira.txt
              python qagen.py INFOSYS-4754 --ppt ./decrypted.pptx --jira-text ./jira.txt
              python qagen.py INFOSYS-4754 --interactive
        """),
    )
    parser.add_argument('jira_num', help='Jira 티켓 번호 (예: INFOSYS-4754)')
    parser.add_argument('--ppt', help='복호화된 PPT 파일 경로')
    parser.add_argument('--jira-text', help='Jira 본문 텍스트 파일 경로')
    parser.add_argument('--interactive', action='store_true', help='직접 텍스트 입력 모드')
    parser.add_argument('--topics', default='',
                        help='QA 주제(테스트 범위). 줄바꿈은 \\n 또는 쉼표로 구분. 지정 시 해당 주제만 생성')
    parser.add_argument('--system-type', choices=['BO', 'SHOP', 'OSS'], default='BO',
                        help='시스템 타입 (기본값: BO)')
    parser.add_argument('--output-dir', default=None,
                        help='출력 디렉토리 (기본값: outputs/)')

    args = parser.parse_args()

    # 입력 소스 확인
    if not args.ppt and not args.jira_text and not args.interactive:
        parser.error("--ppt, --jira-text, --interactive 중 하나 이상 지정해야 합니다.")

    # .env 로드
    load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), '.env'))
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        print("[오류] ANTHROPIC_API_KEY가 설정되지 않았습니다.", file=sys.stderr)
        print("  .env 파일에 ANTHROPIC_API_KEY=sk-ant-... 형태로 입력하세요.", file=sys.stderr)
        sys.exit(1)

    # 입력 텍스트 수집
    input_parts = []

    if args.ppt:
        print(f"[*] PPT 텍스트 추출 중: {args.ppt}")
        ppt_text = extract_ppt_text(args.ppt)
        input_parts.append(f"## PPT 내용\n{ppt_text}")
        print(f"[*] PPT 추출 완료 ({len(ppt_text)} 글자)")

    if args.jira_text:
        print(f"[*] Jira 텍스트 읽는 중: {args.jira_text}")
        jira_text = read_text_file(args.jira_text)
        input_parts.append(f"## Jira 티켓 내용\n{jira_text}")
        print(f"[*] Jira 텍스트 로드 완료 ({len(jira_text)} 글자)")

    if args.interactive:
        interactive_text = interactive_input()
        input_parts.append(f"## 사용자 입력\n{interactive_text}")

    input_text = "\n\n".join(input_parts)

    # QA 주제 정규화 (쉼표/\n 입력을 한 줄에 하나씩으로)
    qa_topics = args.topics.replace('\\n', '\n').replace(',', '\n').strip() if args.topics else ''

    # Claude API 호출
    result = call_claude_api(args.jira_num, input_text, api_key, qa_topics=qa_topics)

    title = result.get('title', args.jira_num)
    system_type = result.get('system_type', args.system_type)
    test_cases = result['test_cases']

    print(f"[*] 생성된 테스트 케이스: {len(test_cases)}건")
    print(f"[*] 문서 제목: {title}")
    print(f"[*] 시스템 타입: {system_type}")

    # xlsx 생성
    wb, file_name = build_workbook(args.jira_num, title, test_cases, system_type=system_type)

    out_dir = args.output_dir or os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')
    out_path = save_workbook(wb, file_name, out_dir)

    print(f"\n[완료] 파일 생성됨: {out_path}")
    print(f"[완료] 크기: {os.path.getsize(out_path)} bytes")
    print(f"[완료] 테스트 케이스: {len(test_cases)}건")


if __name__ == '__main__':
    main()
