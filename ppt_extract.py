#!/usr/bin/env python3
"""
PPT 텍스트 추출기

복호화된 PPT 파일에서 슬라이드별 텍스트를 추출하여 inputs/ 에 저장한다.
API 호출 없음.

사용법:
  python ppt_extract.py path/to/decrypted.pptx
  python ppt_extract.py path/to/decrypted.pptx --output inputs/my_output.txt
"""
import argparse
import os
import sys


def extract_ppt_text(ppt_path):
    """python-pptx로 PPT에서 슬라이드별 텍스트를 추출한다.

    추출 대상: 텍스트 상자, 표 셀, 그룹 도형 내부 텍스트, 발표자 노트.
    (이미지·화면 캡처 속 글자는 OCR하지 않으므로 추출되지 않는다.)
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def collect_shape_text(shape, texts):
        # 그룹 도형은 내부 도형까지 재귀로 들어가 추출한다.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                collect_shape_text(sub, texts)
            return
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
        if shape.has_table:
            for row in shape.table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))

    prs = Presentation(ppt_path)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            collect_shape_text(shape, texts)

        # 발표자 노트
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[발표자 노트] {notes}")

        if texts:
            slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(texts))

    return "\n\n".join(slides_text)


def main():
    parser = argparse.ArgumentParser(description="PPT 텍스트 추출기")
    parser.add_argument('ppt_path', help='복호화된 PPT 파일 경로')
    parser.add_argument('--output', '-o', default=None,
                        help='출력 파일 경로 (기본: inputs/{PPT파일명}.txt)')
    args = parser.parse_args()

    if not os.path.isfile(args.ppt_path):
        print(f"[오류] 파일을 찾을 수 없습니다: {args.ppt_path}", file=sys.stderr)
        sys.exit(1)

    print(f"[*] PPT 텍스트 추출 중: {args.ppt_path}")
    text = extract_ppt_text(args.ppt_path)

    if not text.strip():
        print("[경고] 추출된 텍스트가 없습니다.", file=sys.stderr)
        sys.exit(1)

    # 출력 경로 결정
    if args.output:
        out_path = args.output
    else:
        base_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'inputs')
        os.makedirs(base_dir, exist_ok=True)
        base_name = os.path.splitext(os.path.basename(args.ppt_path))[0]
        out_path = os.path.join(base_dir, f"{base_name}.txt")

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(text)

    print(f"[완료] 저장됨: {out_path}")
    print(f"[완료] {len(text)} 글자, 슬라이드 {text.count('[슬라이드')}장")
    print()
    print("다음 단계: Claude Code 세션에서 아래와 같이 요청하세요.")
    print(f'  "inputs/{os.path.basename(out_path)} 읽고 JIRA번호 XXXX-0000 으로 QA 문서 생성해줘"')


if __name__ == '__main__':
    main()
