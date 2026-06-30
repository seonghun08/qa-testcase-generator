"""
QA 문서 자동 생성 - 로컬 웹 앱

Claude Code CLI (Pro 구독)를 활용하여 API 비용 없이 QA 테스트 케이스 문서를 생성한다.

실행: streamlit run app.py
"""
import json
import os
import subprocess
import sys
import textwrap

import pandas as pd
import streamlit as st

# 프로젝트 루트를 path에 추가
PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, PROJECT_DIR)

from template_builder import build_workbook, save_workbook

SETTINGS_PATH = os.path.join(PROJECT_DIR, "settings.json")
RULES_PATH = os.path.join(PROJECT_DIR, "rules.json")


def _rel(path):
    """표시용 경로: 설치 위치(PROJECT_DIR) 기준 상대경로로 변환한다.

    내부 경로는 항상 PROJECT_DIR 기준으로 동적 계산되므로 PC마다 달라도 정상 동작한다.
    화면에는 절대경로 대신 상대경로(예: outputs/파일.xlsx)만 보여 노이즈를 줄인다.
    """
    if not path:
        return ""
    try:
        return os.path.relpath(path, PROJECT_DIR).replace("\\", "/")
    except ValueError:
        return os.path.basename(path)


# ============================================================
# 설정 파일 관리
# ============================================================

def load_settings():
    """settings.json에서 설정값을 읽는다."""
    defaults = {
        "developer_name": "",
        "tester_name": "",
        "qa_name": "",
        "default_system_type": "BO",
    }
    if os.path.isfile(SETTINGS_PATH):
        with open(SETTINGS_PATH, "r", encoding="utf-8") as f:
            saved = json.load(f)
        defaults.update(saved)
    return defaults


def save_settings(settings):
    """settings.json에 설정값을 저장한다."""
    with open(SETTINGS_PATH, "w", encoding="utf-8") as f:
        json.dump(settings, f, ensure_ascii=False, indent=2)


def load_rules():
    """rules.json에서 커스텀 규칙을 읽는다.

    사내 DRM이 .txt를 암호화해 읽지 못하는 문제 때문에 .json으로 관리한다.
    형식: {"rules": ["규칙1", "규칙2", ...]}
    """
    data = load_rules_data()
    rules = [str(r).strip() for r in data.get("rules", []) if str(r).strip()]
    return "\n".join(rules)


def load_rules_data():
    """rules.json 전체(딕셔너리)를 읽는다. _comment 보존용."""
    if not os.path.isfile(RULES_PATH):
        return {"rules": []}
    try:
        with open(RULES_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return {"rules": []}


def save_rules_list(rules_list, comment=None):
    """규칙 리스트를 rules.json에 저장한다. _comment는 보존한다."""
    data = {}
    if comment:
        data["_comment"] = comment
    data["rules"] = [str(r).strip() for r in rules_list if str(r).strip()]
    with open(RULES_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


# ============================================================
# PPT 텍스트 추출
# ============================================================

def extract_ppt_text(ppt_bytes):
    """업로드된 PPT 바이트에서 텍스트를 추출한다.

    추출 대상: 텍스트 상자, 표 셀, 그룹 도형 내부 텍스트, 발표자 노트.
    (이미지·화면 캡처 속 글자는 OCR하지 않으므로 추출되지 않는다.)
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import io

    def collect_shape_text(shape, texts):
        # 그룹 도형은 내부 도형까지 재귀로 들어가 추출한다.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                collect_shape_text(sub, texts)
            return
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
        if shape.has_table:
            for row in shape.table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))

    prs = Presentation(io.BytesIO(ppt_bytes))
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            collect_shape_text(shape, texts)

        # 발표자 노트
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[발표자 노트] {notes}")

        if texts:
            slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(texts))

    return "\n\n".join(slides_text)


# ============================================================
# Claude Code CLI 호출
# ============================================================

CLAUDE_PROMPT_TEMPLATE = textwrap.dedent("""\
    당신은 PG사 백오피스 QA 테스트 케이스 작성 전문가입니다.
    아래 개발 요청서를 분석하여 사내 표준 QA 테스트 케이스를 JSON으로 생성하세요.

    ## 출력 형식
    반드시 아래 JSON 형식으로만 응답하세요. 설명이나 마크다운 코드블록 없이 순수 JSON만 출력하세요.

    {{
      "title": "문서 제목",
      "system_type": "BO 또는 SHOP 또는 OSS (문서 기본값)",
      "test_cases": [
        {{
          "id": "TB_01_001",
          "system_type": "BO|SHOP|OSS 중 해당하는 것. 같은 메뉴가 둘 다에 적용되면 \\"BO/SHOP\\"처럼 슬래시로 결합",
          "precond": "사전조건 (예: 내부직원 계정 로그인\\n\\nPG 정산 관리 > 메뉴경로)",
          "step": "조회|등록|수정|삭제",
          "category": "구분",
          "item": "항목",
          "procedure": "수행절차 (번호 매기기)",
          "expected": "기대결과 (번호 매기기)",
          "log_db_type": "select|insert|update"
        }}
      ]
    }}

    ## TestCaseID 규칙
    - 형식: TB_{{기능그룹번호}}_{{순번}} (예: TB_01_001)
    - 화면 단위로 그룹 분리 (메인 화면=01, 상세 팝업=02 등)
    - 한 화면 안에서는 조회→등록→수정→삭제 순
    - 순번은 그룹마다 001부터 시작

    ## 시스템 타입 지정 (각 케이스의 system_type)
    - 케이스마다 어느 시스템에서 테스트하는지 지정한다: BO, SHOP, OSS 중 하나.
    - 같은 메뉴/기능이 BO와 SHOP 양쪽에 동일하게 적용되면 "BO/SHOP"으로 결합해 한 케이스로 작성한다(케이스를 따로 쪼개지 않음).
    - OSS는 보통 별도 시스템이다. 명확히 BO/SHOP과 같은 화면을 공유한다는 근거가 없으면 OSS를 BO/SHOP과 결합하지 말고 단독("OSS")으로 둔다.
    - 판단이 모호하면 문서 기본값({system_default})을 사용한다.

    ## 사전조건 형식
    - "{{계정유형}} 계정 로그인\\n\\n{{메뉴경로}}"
    - BO→내부직원, SHOP→가맹점관리자, OSS→운영지원담당자

    ## QA 관점 (가장 중요 — 모든 케이스에 우선 적용)
    - 이 QA는 **비개발자(기획·운영 담당자)** 가 화면을 직접 클릭하며 수행한다. 목적은 "UI/UX가 정상적으로 동작하는가"를 눈으로 확인하는 것이다.
    - 수행절차·기대결과에는 **화면에서 눈으로 확인 가능한 것만** 쓴다:
      버튼/입력/필터/검색 동작, 목록·팝업·메시지·알럿 노출 여부와 문구, 정렬·페이징, 검색 결과가 입력 조건에 맞게 나오는지, 권한/시스템별 화면 차이(노출·미노출).
    - **다음은 비개발자가 확인할 수 없으므로 수행절차·기대결과에서 제외한다**(개발자가 로그/DB 섹션에서 확인할 영역):
      암호화·해시 처리 여부, 내부 매칭 로직(완전일치/범위검색 등 알고리즘 설명), DB 테이블/컬럼명, 적재 테이블, API명, 쿼리, 트랜잭션 처리 등 화면 밖 내부 구현.
      → 검색은 "조건을 입력하면 결과가 그 조건에 맞게 표시되는가" 수준으로만 기술한다.

    ## 테스트 케이스 작성 지침
    - 분량: 케이스 개수를 미리 정하지 않는다. 'QA 주제'가 지정되면 그 주제에만 한정해 주제당 1건(정상·에러 흐름이 많으면 최대 2건)으로 작성하고, 나열되지 않은 기능/화면은 만들지 않는다.
      - 'QA 주제'가 없으면 개발 요청서에서 핵심 기능을 뽑아 화면/기능 단위로 크게 묶어 간결하게 작성한다(사람이 직접 QA하므로 세부 단계까지 쪼개지 않음).
      - 한 케이스의 수행절차·기대결과는 핵심 확인 포인트만 간결하게 쓴다. 세부 동작을 단계마다 잘게 쪼개 모두 나열하지 말고, 비개발자가 화면에서 자유롭게 조작하며 확인할 수 있는 추상 수준으로 묶어서 표현한다(검증 포인트가 많아도 핵심만 남김).
      - 확인 팝업/필수 입력 검증/정상 처리/실패 처리/이력 적재는 별도 케이스로 분리하지 말고 한 케이스에 통합한다.
    - **목록 조회와 상세 조회는 가능하면 별도 케이스(별도 기능 그룹)로 분리한다.**
      - 목록 조회: 검색/필터 입력, 목록 노출, 정렬·페이징, 요약(건수/금액) 표시 확인.
      - 상세 조회: 목록 항목 클릭 시 상세 팝업/화면이 열리고 표시 항목이 올바르게 노출되는지 확인.
    - 수행절차: 화면에서 따라 할 수 있는 조작 단계 (번호 형식, 보통 3~5단계 내외로 간결하게)
    - 기대결과: 각 단계에 대응하는, 화면에서 확인 가능한 결과 (핵심만; 자명한 결과는 생략)
    - 정상 케이스 + 사용자가 화면에서 마주칠 수 있는 에러/경계 케이스(필수 미입력, 형식 오류 메시지 등) 포함
    - log_db_type: select(조회), insert(등록), update(수정/삭제/상태변경)

    {custom_rules}

    {qa_topics}

    ## Jira 번호: {jira_num}

    ## 개발 요청 내용:
    {input_text}

    위 내용을 기반으로 JSON을 생성하세요. JSON만 출력하세요.
""")


def find_claude_cli():
    """Claude Code CLI 실행 경로를 찾는다."""
    import shutil

    for name in ["claude.cmd", "claude.exe", "claude"]:
        path = shutil.which(name)
        if path:
            return path

    candidates = [
        os.path.join(os.environ.get("APPDATA", ""), "npm", "claude.cmd"),
        r"C:\npm-global\claude.cmd",
        os.path.expanduser("~\\AppData\\Roaming\\npm\\claude.cmd"),
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p

    return None


def call_claude_cli(jira_num, input_text, custom_rules="", qa_topics="", system_default="BO"):
    """Claude Code CLI를 호출하여 테스트 케이스 JSON을 받는다."""
    claude_path = find_claude_cli()
    if not claude_path:
        raise FileNotFoundError("claude CLI를 찾을 수 없습니다.")

    # 커스텀 규칙을 프롬프트 섹션으로 포맷
    rules_section = ""
    if custom_rules.strip():
        rules_section = f"## 추가 규칙 (반드시 준수)\n{custom_rules}"

    # QA 주제(테스트 범위)를 프롬프트 섹션으로 포맷
    topics_section = ""
    if qa_topics.strip():
        topics_section = (
            "## QA 주제 (사용자 지정 테스트 범위 — 반드시 이 주제들에만 한정)\n"
            f"{qa_topics}\n\n"
            "- 위 주제만 테스트 케이스로 작성하고, 주제당 1건(정상·에러 흐름이 많으면 최대 2건)으로 통합한다.\n"
            "- 나열되지 않은 기능/화면은 절대 생성하지 않는다."
        )

    prompt = CLAUDE_PROMPT_TEMPLATE.format(
        jira_num=jira_num,
        input_text=input_text,
        custom_rules=rules_section,
        qa_topics=topics_section,
        system_default=system_default,
    )

    # 사내 DRM이 .txt를 암호화하는 문제 때문에 임시 파일을 만들지 않고,
    # 프롬프트를 stdin으로 직접 전달한다.
    cmd = f'"{claude_path}" -p --output-format text'

    result = subprocess.run(
        cmd,
        input=prompt,
        capture_output=True,
        text=True,
        timeout=300,
        cwd=PROJECT_DIR,
        shell=True,
        encoding="utf-8",
        errors="replace",
    )

    stdout = result.stdout or ""
    stderr = result.stderr or ""

    if stdout.strip():
        return stdout.strip()

    raise RuntimeError(
        f"Claude CLI가 빈 응답을 반환했습니다. (exit {result.returncode})\n"
        f"stderr: {stderr[:500]}"
    )


def parse_claude_response(response_text):
    """Claude 응답에서 JSON을 파싱한다."""
    text = response_text.strip()

    if "```json" in text:
        start = text.index("```json") + len("```json")
        end = text.index("```", start)
        text = text[start:end].strip()
    elif "```" in text:
        start = text.index("```") + 3
        end = text.index("```", start)
        text = text[start:end].strip()

    if "{" in text:
        start = text.index("{")
        end = text.rindex("}") + 1
        text = text[start:end]

    return json.loads(text)


# ============================================================
# Streamlit UI
# ============================================================

def render_rules_editor():
    """rules.json 규칙을 표로 편집/저장하는 UI (다음 생성부터 적용)."""
    rules_data = load_rules_data()
    rules_list = [str(r).strip() for r in rules_data.get("rules", []) if str(r).strip()]

    st.caption(
        f"현재 {len(rules_list)}개 · 표에서 직접 **수정**, 맨 아래 빈 행에 입력해 **추가**, "
        "행 선택 후 삭제할 수 있습니다. 셀을 더블클릭하면 전체 내용이 보입니다."
    )

    df = pd.DataFrame({"규칙": rules_list})
    edited = st.data_editor(
        df,
        num_rows="dynamic",
        use_container_width=True,
        hide_index=True,
        column_config={
            "규칙": st.column_config.TextColumn(
                "규칙 (한 행에 하나)", width="large",
                help="QA 생성 시 AI가 참고하는 규칙입니다.",
            ),
        },
        key="rules_editor",
    )

    col_a, col_b, _ = st.columns([1, 1, 3])
    if col_a.button("💾 규칙 저장", type="primary", use_container_width=True):
        new_rules = [str(r).strip() for r in edited["규칙"].tolist() if str(r).strip()]
        save_rules_list(new_rules, comment=rules_data.get("_comment"))
        st.success(f"{len(new_rules)}개 규칙 저장됨 — 다음 생성부터 적용됩니다.")
    if col_b.button("↩️ 되돌리기", use_container_width=True,
                    help="저장 전 편집 내용을 취소하고 파일에 저장된 상태로 복원"):
        st.session_state.pop("rules_editor", None)
        st.rerun()

    st.caption(f"📂 `{_rel(RULES_PATH)}`")


def main():
    st.set_page_config(page_title="QA 문서 자동 생성", page_icon="📋", layout="wide")

    # 화면 배율(90%) + 본문을 화면 전체 폭으로 사용 (좌우 여백 최소화)
    st.markdown(
        "<style>"
        ".stApp { zoom: 0.9; }"
        ".block-container { max-width: 100% !important; padding-left: 1.5rem; padding-right: 1.5rem; }"
        "[data-testid=\"stAppDeployButton\"] { display: none; }"  # 외부 배포용 버튼 숨김(사내 로컬 사용)
        "</style>",
        unsafe_allow_html=True,
    )

    st.title("📋 QA 테스트 케이스 문서 자동 생성")
    st.caption("Claude Code (Pro 구독) 기반")

    # 설정 로드
    settings = load_settings()
    custom_rules = load_rules()

    # --- 사이드바: 설정 ---
    with st.sidebar:
        st.subheader("👤 담당자 설정")
        developer_name = st.text_input("개발자명", value=settings.get("developer_name", ""),
                                       placeholder="홍길동")
        tester_name = st.text_input("테스터명", value=settings.get("tester_name", ""),
                                    placeholder="(선택)")
        qa_name = st.text_input("QA명", value=settings.get("qa_name", ""),
                                placeholder="(선택)")
        if st.button("💾 설정 저장", use_container_width=True):
            save_settings({
                "developer_name": developer_name,
                "tester_name": tester_name,
                "qa_name": qa_name,
            })
            st.success("설정 저장됨!")

        st.divider()
        st.markdown("""
        **사용법**
        1. 담당자 설정 입력 (필요 시 저장)
        2. Jira 티켓 입력
        3. PPT 업로드 또는 텍스트 입력
        4. 문서 생성 클릭
        5. 미리보기 결과 편집 (선택)
        6. xlsx 다운로드

        **규칙 수정**
        **📏 작성 규칙 편집** 탭에서 추가/수정/삭제하면
        다음 생성부터 적용됩니다.
        """)

    # --- 메인 영역: 탭 구성 ---
    tab_gen, tab_hist, tab_rules, tab_env = st.tabs([
        "📋 QA 문서 생성",
        "📂 생성 이력",
        "📏 작성 규칙 편집",
        "⚙️ 환경 점검 / 설치 가이드",
    ])

    with tab_gen:
        render_generation_tab(developer_name, tester_name, qa_name, custom_rules)

    with tab_hist:
        render_history_tab(developer_name, tester_name, qa_name)

    with tab_env:
        st.subheader("⚙️ 환경 점검 / 설치 가이드")
        st.caption("앱 실행·문서 생성에 필요한 구성요소를 점검하고, 빠진 항목은 설치 방법을 안내합니다.")
        if st.button("환경 점검 실행"):
            st.session_state["_run_setup_check"] = True
        if st.session_state.get("_run_setup_check"):
            from setup_guide import render_setup_check
            render_setup_check()

    with tab_rules:
        render_rules_editor()


def render_generation_tab(developer_name, tester_name, qa_name, custom_rules):
    """QA 문서 생성 탭: 입력 → Claude 호출 → xlsx 생성/다운로드."""
    jira_num = st.text_input(
        "Jira 티켓 번호",
        placeholder="INFOSYS-0000",
        help="생성 문서의 파일명과 내부 제목에 사용됩니다.",
    )

    # --- 입력 (PPT 업로드 + 자유 입력 한 칸으로 통합) ---
    uploaded_file = st.file_uploader(
        "📎 복호화된 PPT 파일 (선택)",
        type=["pptx"],
    )
    ppt_text = ""
    if uploaded_file:
        with st.spinner("PPT 텍스트 추출 중..."):
            ppt_text = extract_ppt_text(uploaded_file.read())
        st.success(f"추출 완료: {len(ppt_text)} 글자")
        with st.expander("추출된 텍스트 미리보기"):
            st.text(ppt_text[:2000] + ("..." if len(ppt_text) > 2000 else ""))

    manual_text = st.text_area(
        "📝 요청 내용 · QA 작성 지침 (자유 입력)",
        height=400,
        key="main_input",
        placeholder=(
            "개발 요청 내용, 신규/변경 메뉴, 화면 동작, QA 시 확인할 점 등을 자유롭게 적으세요.\n"
            "PPT를 올렸다면 PPT에 없는 추가 설명만 적어도 됩니다."
        ),
    )

    # --- 입력 텍스트 결합 ---
    input_parts = []
    if ppt_text:
        input_parts.append(f"## PPT 내용\n{ppt_text}")
    if manual_text:
        input_parts.append(f"## 요청 내용 및 QA 지침\n{manual_text}")

    input_text = "\n\n".join(input_parts)

    # QA 주제 입력칸은 제거됨 — 위 자유 입력 한 칸에 통합한다.
    qa_topics = ""

    # --- 생성 버튼 ---
    st.divider()

    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        generate_btn = st.button(
            "🚀 문서 생성",
            type="primary",
            use_container_width=True,
            disabled=not (jira_num and input_text),
        )

    if not jira_num and input_text:
        st.warning("Jira 티켓 번호를 입력하세요.")

    if generate_btn and jira_num and input_text:
        with st.status("QA 테스트 케이스 생성 중...", expanded=True) as status:
            st.write("🤖 Claude Code CLI 호출 중... (최대 5분)")

            try:
                raw_response = call_claude_cli(jira_num, input_text, custom_rules=custom_rules, qa_topics=qa_topics)
            except subprocess.TimeoutExpired:
                st.error("⏱️ Claude CLI 응답 시간 초과 (5분). 입력을 줄이고 다시 시도하세요.")
                st.stop()
            except FileNotFoundError as e:
                st.error(f"❌ {e}")
                st.stop()
            except RuntimeError as e:
                st.error(f"❌ {e}")
                st.stop()

            st.write("📋 응답 파싱 중...")

            try:
                data = parse_claude_response(raw_response)
            except (json.JSONDecodeError, ValueError) as e:
                st.error(f"❌ JSON 파싱 실패: {e}")
                with st.expander("원본 응답 보기"):
                    st.code(raw_response[:3000])
                st.stop()

            title = data.get("title", jira_num)
            # 문서 기본 시스템 타입 (케이스가 별도 지정 안 했을 때 fallback). 케이스별 값이 우선.
            resp_system_type = data.get("system_type", "BO")
            test_cases = data.get("test_cases", [])

            if not test_cases:
                st.error("❌ 생성된 테스트 케이스가 없습니다.")
                with st.expander("원본 응답 보기"):
                    st.code(raw_response[:3000])
                st.stop()

            st.write(f"📊 테스트 케이스 {len(test_cases)}건 생성 완료")
            st.write("📄 xlsx 파일 생성 중...")
            built = build_and_save(jira_num, title, resp_system_type, test_cases,
                                   developer_name, tester_name, qa_name)
            status.update(label="✅ 생성 완료!", state="complete")

        # 생성 결과를 세션에 보존 (편집/재생성 시 rerun 되어도 유지)
        st.session_state["gen_result"] = {
            "jira_num": jira_num,
            "title": title,
            "system_type": resp_system_type,
            "test_cases": test_cases,
        }
        st.session_state["last_build"] = built
        st.session_state.pop("cases_editor", None)  # 새 생성이면 편집기 초기화
        st.success(f"**{built['name']}** 생성 완료! ({built['count']}건) — 아래 표에서 편집 후 재생성할 수 있습니다.")

    # --- 생성 결과 편집 / 다운로드 (세션에 결과가 있으면 항상 표시) ---
    if st.session_state.get("gen_result"):
        render_result_editor(developer_name, tester_name, qa_name)


def build_and_save(jira_num, title, system_type, test_cases,
                   developer_name, tester_name, qa_name):
    """test_cases로 xlsx + JSON을 outputs/에 저장하고 경로 정보를 반환한다."""
    wb, file_name = build_workbook(
        jira_num, title, test_cases,
        system_type=system_type,
        developer_name=developer_name,
        tester_name=tester_name,
        qa_name=qa_name,
    )
    out_dir = os.path.join(PROJECT_DIR, "outputs")
    out_path = save_workbook(wb, file_name, out_dir)
    json_path = os.path.join(out_dir, f"{jira_num}_test_cases.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({"title": title, "system_type": system_type, "test_cases": test_cases},
                  f, ensure_ascii=False, indent=2)
    return {"path": out_path, "name": file_name, "json_path": json_path,
            "count": len(test_cases)}


def render_result_editor(developer_name, tester_name, qa_name,
                         result_key="gen_result", editor_key="cases_editor",
                         build_key="last_build", show_close=False, close_flag=None):
    """케이스를 표로 인라인 편집하고, 편집본으로 xlsx를 재생성한다.

    result_key/editor_key/build_key를 바꾸면 다른 탭에서도 독립적으로 재사용된다.
    show_close=True이면 헤더 우측 상단에 '편집 닫기' 버튼을 표시한다.
    close_flag를 주면 닫기 시 그 플래그만 세팅하고, 실제 상태 초기화는 호출부가
    (위젯 생성 전에) 처리한다 — 선택 위젯까지 안전하게 리셋하기 위함.
    """
    res = st.session_state[result_key]
    st.divider()
    if show_close:
        hcol, bcol = st.columns([5, 1], vertical_alignment="center")
        hcol.subheader(f"📝 결과 편집 — {res['title']}")
        if bcol.button("편집 닫기", key=f"{editor_key}_close", use_container_width=True):
            if close_flag:
                st.session_state[close_flag] = True
            else:
                for k in (result_key, editor_key, build_key):
                    st.session_state.pop(k, None)
            st.rerun()
    else:
        st.subheader(f"📝 결과 편집 — {res['title']}")
    st.caption(
        "표에서 직접 수정 / 맨 아래 빈 행에 추가 / 행 선택 후 삭제할 수 있습니다. "
        "긴 셀(수행절차·기대결과)은 더블클릭하면 전체가 보입니다. "
        "편집 후 '편집 반영 + xlsx 재생성'을 누르세요."
    )

    cols = ["id", "system_type", "step", "category", "item",
            "procedure", "expected", "precond", "log_db_type"]
    df = pd.DataFrame([{c: tc.get(c, "") for c in cols} for tc in res["test_cases"]])

    edited = st.data_editor(
        df,
        num_rows="dynamic",
        use_container_width=True,
        hide_index=True,
        column_config={
            "id": st.column_config.TextColumn("ID", width="small"),
            "system_type": st.column_config.TextColumn("시스템", width="small",
                help="BO / SHOP / OSS (또는 BO/SHOP)"),
            "step": st.column_config.TextColumn("단계", width="small"),
            "category": st.column_config.TextColumn("구분", width="medium"),
            "item": st.column_config.TextColumn("항목", width="medium"),
            "procedure": st.column_config.TextColumn("수행절차", width="large"),
            "expected": st.column_config.TextColumn("기대결과", width="large"),
            "precond": st.column_config.TextColumn("사전조건", width="medium"),
            "log_db_type": st.column_config.TextColumn("로그/DB", width="small",
                help="select / insert / update"),
        },
        key=editor_key,
    )

    c1, c2, _ = st.columns([1, 1, 2])
    if c1.button("💾 편집 반영 + xlsx 재생성", type="primary", use_container_width=True,
                 key=f"{editor_key}_save"):
        new_cases = []
        for r in edited.to_dict("records"):
            rec = {k: ("" if pd.isna(v) else v) for k, v in r.items()}
            if str(rec.get("item", "")).strip() or str(rec.get("procedure", "")).strip():
                new_cases.append(rec)
        if not new_cases:
            st.warning("케이스가 비어 있습니다. 한 건 이상 입력하세요.")
        else:
            built = build_and_save(res["jira_num"], res["title"], res["system_type"],
                                   new_cases, developer_name, tester_name, qa_name)
            st.session_state[build_key] = built
            st.success(f"편집본으로 재생성 완료! ({built['count']}건)")

    if c2.button("↩️ 편집 취소", use_container_width=True, key=f"{editor_key}_reset",
                 help="편집 내용을 버리고 원본 생성 결과로 되돌립니다."):
        st.session_state.pop(editor_key, None)
        st.rerun()

    # 다운로드 (가장 최근 빌드본)
    built = st.session_state.get(build_key)
    if built and os.path.isfile(built["path"]):
        with open(built["path"], "rb") as f:
            st.download_button(
                label=f"📥 {built['name']} 다운로드 ({built['count']}건)",
                data=f.read(),
                file_name=built["name"],
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                type="primary",
                key=f"{build_key}_dl",
            )
        st.caption(f"파일: `{_rel(built['path'])}`")
        with st.expander("생성된 JSON 보기"):
            try:
                with open(built["json_path"], "r", encoding="utf-8") as f:
                    st.json(json.load(f))
            except OSError:
                pass


def list_history():
    """outputs/의 생성 이력(*_test_cases.json 기준)을 최신순으로 반환한다."""
    out_dir = os.path.join(PROJECT_DIR, "outputs")
    if not os.path.isdir(out_dir):
        return []
    all_files = os.listdir(out_dir)
    xlsx_files = [f for f in all_files if f.endswith(".xlsx")]
    items = []
    for name in all_files:
        if not name.endswith("_test_cases.json") or ".bak" in name:
            continue
        path = os.path.join(out_dir, name)
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except (OSError, json.JSONDecodeError):
            continue
        jira_num = name[:-len("_test_cases.json")]
        xlsx = next((os.path.join(out_dir, f) for f in xlsx_files
                     if f.startswith(f"TB_TC_{jira_num}_")), None)
        items.append({
            "jira_num": jira_num,
            "title": data.get("title", jira_num),
            "system_type": data.get("system_type", "BO"),
            "count": len(data.get("test_cases", [])),
            "json_path": path,
            "xlsx_path": xlsx,
            "mtime": os.path.getmtime(path),
            "data": data,
        })
    items.sort(key=lambda x: x["mtime"], reverse=True)
    return items


def render_history_tab(developer_name, tester_name, qa_name):
    """생성 이력 조회 / 재다운로드 / 이 탭에서 바로 편집·재생성 / 삭제."""
    from datetime import datetime

    # 편집 닫기 처리: 위젯 생성 전에 관련 상태(선택 항목 포함)를 모두 초기화한다.
    if st.session_state.pop("_hist_close", False):
        for k in ("hist_result", "hist_cases_editor", "hist_last_build", "hist_sel"):
            st.session_state.pop(k, None)

    items = list_history()
    total = len(items)

    if total == 0:
        st.info("아직 생성 이력이 없습니다. '📋 QA 문서 생성' 탭에서 먼저 문서를 만들어 보세요.")
    else:
        PAGE_SIZE = 10
        pages = (total + PAGE_SIZE - 1) // PAGE_SIZE
        page = min(max(st.session_state.get("hist_page", 1), 1), pages)
        start = (page - 1) * PAGE_SIZE
        page_items = items[start:start + PAGE_SIZE]

        # 목록(표)
        rows = [{
            "No": start + i + 1,
            "Jira 번호": it["jira_num"],
            "제목": it["title"],
            "건수": it["count"],
            "시스템": it["system_type"],
            "생성일시": datetime.fromtimestamp(it["mtime"]).strftime("%Y-%m-%d %H:%M"),
            "xlsx": "있음" if it["xlsx_path"] else "없음",
        } for i, it in enumerate(page_items)]
        st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

        # 페이지 이동
        pc1, pc2, pc3 = st.columns([1, 1, 6])
        if pc1.button("◀ 이전", disabled=(page <= 1), use_container_width=True):
            st.session_state["hist_page"] = page - 1
            st.rerun()
        if pc2.button("다음 ▶", disabled=(page >= pages), use_container_width=True):
            st.session_state["hist_page"] = page + 1
            st.rerun()
        pc3.markdown(f"**{page} / {pages}** 페이지 · 총 {total}건")

        # 작업 대상 선택 + 액션
        st.divider()
        labels = {f"{it['jira_num']} · {it['title']} ({it['count']}건)": it for it in page_items}
        sel_label = st.selectbox("작업할 항목 선택 (현재 페이지)", list(labels.keys()), key="hist_sel")
        sel = labels[sel_label]

        a1, a2, a3, _ = st.columns([1.4, 1.6, 1, 3])
        if sel["xlsx_path"] and os.path.isfile(sel["xlsx_path"]):
            with open(sel["xlsx_path"], "rb") as f:
                a1.download_button(
                    "📥 xlsx 다운로드", data=f.read(),
                    file_name=os.path.basename(sel["xlsx_path"]),
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key="hist_dl", use_container_width=True,
                )
        else:
            a1.button("📥 xlsx 없음", disabled=True, use_container_width=True)

        if a2.button("✏️ 편집기로 불러오기", key="hist_load", type="primary", use_container_width=True):
            st.session_state["hist_result"] = {
                "jira_num": sel["jira_num"],
                "title": sel["title"],
                "system_type": sel["system_type"],
                "test_cases": sel["data"].get("test_cases", []),
            }
            if sel["xlsx_path"] and os.path.isfile(sel["xlsx_path"]):
                st.session_state["hist_last_build"] = {
                    "path": sel["xlsx_path"],
                    "name": os.path.basename(sel["xlsx_path"]),
                    "json_path": sel["json_path"],
                    "count": sel["count"],
                }
            else:
                st.session_state.pop("hist_last_build", None)
            st.session_state.pop("hist_cases_editor", None)
            st.rerun()

        if a3.button("🗑 삭제", key="hist_del", use_container_width=True):
            st.session_state["_pending_delete"] = sel["jira_num"]

        if st.session_state.get("_pending_delete") == sel["jira_num"]:
            st.warning(f"'{sel['jira_num']}' 관련 파일(JSON · xlsx)을 삭제할까요?")
            d1, d2, _ = st.columns([1, 1, 4])
            if d1.button("삭제 확인", key="hist_delok", type="primary"):
                try:
                    os.remove(sel["json_path"])
                    if sel["xlsx_path"] and os.path.isfile(sel["xlsx_path"]):
                        os.remove(sel["xlsx_path"])
                except OSError as e:
                    st.error(f"삭제 실패: {e}")
                st.session_state.pop("_pending_delete", None)
                st.rerun()
            if d2.button("취소", key="hist_delno"):
                st.session_state.pop("_pending_delete", None)
                st.rerun()

    # --- 선택한 이력을 이 탭에서 바로 편집/재생성 ---
    if st.session_state.get("hist_result"):
        render_result_editor(developer_name, tester_name, qa_name,
                             result_key="hist_result", editor_key="hist_cases_editor",
                             build_key="hist_last_build", show_close=True, close_flag="_hist_close")


if __name__ == "__main__":
    main()
